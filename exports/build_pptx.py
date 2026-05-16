"""
Genera la presentación ejecutiva .pptx para defensa frente al panel MELI.

Estructura: Minto Pyramid Principle (conclusión-primero) + SCQA + 1 idea/slide.
Tono: dark template institucional MELI, Tufte data-ink alto, ~28pt+ font, máx 3-4 líneas/slide.
Target: 20-30 min de defensa frente a panel mixto (HR + Tech + Ops Regional).

Output: business_case_andres_islas_pitch.pptx (~12 slides).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# ============================================================
# MELI brand
# ============================================================
MELI_YELLOW = RGBColor(0xFF, 0xE6, 0x00)
MELI_YELLOW_DARK = RGBColor(0xFC, 0xD4, 0x00)
MELI_YELLOW_SOFT = RGBColor(0xFF, 0xF7, 0xB8)
MELI_BLUE = RGBColor(0x2D, 0x32, 0x77)
MELI_BLUE_DARK = RGBColor(0x1F, 0x25, 0x62)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x5A, 0x5A, 0x6A)
SUCCESS = RGBColor(0x00, 0xA6, 0x50)
DANGER = RGBColor(0xE6, 0x39, 0x46)
INFO = RGBColor(0x34, 0x83, 0xFA)
PAGE_BG = RGBColor(0xFF, 0xFF, 0xFF)

FONT_HEADING = "Calibri"  # institucional, universal
FONT_BODY = "Calibri"
FONT_MONO = "Consolas"

# 16:9 standard widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUTPUT = Path(__file__).parent / "business_case_andres_islas_pitch.pptx"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]


# ============================================================
# Helpers
# ============================================================

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.line.fill.background() if line_color is None else None
    if line_color is not None:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    if fill_color is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    if not text:
        return tb
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return tb


def add_rich_paragraphs(slide, x, y, w, h, paragraphs, *, anchor=MSO_ANCHOR.TOP):
    """paragraphs is a list of dicts: {text, size, bold, color, italic, align, space_before}."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.vertical_anchor = anchor
    for i, par in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = par.get("align", PP_ALIGN.LEFT)
        if par.get("space_before"):
            p.space_before = Pt(par["space_before"])
        run = p.add_run()
        run.text = par["text"]
        run.font.size = Pt(par.get("size", 14))
        run.font.bold = par.get("bold", False)
        run.font.italic = par.get("italic", False)
        run.font.name = par.get("font", FONT_BODY)
        run.font.color.rgb = par.get("color", INK)
    return tb


def add_brand_strip(slide, page_num=None, total=12):
    """Top yellow band + bottom footer with brand info + page number."""
    # Top thin yellow line
    add_rect(slide, Emu(0), Emu(0), SLIDE_W, Inches(0.08), fill_color=MELI_YELLOW)
    # Bottom footer
    add_rect(slide, Emu(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4), fill_color=MELI_BLUE_DARK)
    add_text(slide, Inches(0.4), SLIDE_H - Inches(0.36),
             Inches(8), Inches(0.3),
             "Andrés Islas Bravo · Innovación y Robótica · Mercado Envíos MX",
             size=10, color=WHITE, font=FONT_BODY)
    if page_num is not None:
        add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.36),
                 Inches(0.9), Inches(0.3),
                 f"{page_num}/{total}",
                 size=10, color=MELI_YELLOW, font=FONT_MONO,
                 align=PP_ALIGN.RIGHT)


def slide_title(slide, eyebrow, title, subtitle=None):
    """Standard slide header: small eyebrow + big title + subtitle."""
    add_text(slide, Inches(0.5), Inches(0.35),
             Inches(7), Inches(0.4),
             eyebrow.upper(),
             size=11, bold=True, color=MELI_YELLOW_DARK, font=FONT_MONO)
    add_text(slide, Inches(0.5), Inches(0.7),
             Inches(12.3), Inches(0.9),
             title,
             size=30, bold=True, color=MELI_BLUE, font=FONT_HEADING)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.55),
                 Inches(12.3), Inches(0.5),
                 subtitle,
                 size=15, italic=True, color=MUTED, font=FONT_BODY)
    # Underline
    add_rect(slide, Inches(0.5), Inches(2.0), Inches(1.2), Emu(35000),
             fill_color=MELI_YELLOW)


# ============================================================
# SLIDE 1 — COVER
# ============================================================

s = prs.slides.add_slide(BLANK)
# Background MELI blue
add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill_color=MELI_BLUE)
# Yellow accent corner
add_rect(s, Emu(0), Emu(0), Inches(2), Inches(0.5), fill_color=MELI_YELLOW)
# Logo dot
add_rect(s, Inches(0.6), Inches(1.0), Inches(0.25), Inches(0.25),
         fill_color=MELI_YELLOW)
add_text(s, Inches(1.0), Inches(1.0), Inches(8), Inches(0.3),
         "MERCADO ENVÍOS MX · INNOVACIÓN Y ROBÓTICA",
         size=12, bold=True, color=MELI_YELLOW, font=FONT_MONO)

add_text(s, Inches(0.6), Inches(1.6), Inches(11), Inches(0.5),
         "Work Sample · Andrés Islas Bravo · Mayo 2026",
         size=14, color=WHITE, font=FONT_BODY)

# Main title
add_rich_paragraphs(s, Inches(0.6), Inches(2.5), Inches(12), Inches(2.5), [
    {"text": "Subir la vara: ", "size": 56, "bold": True, "color": WHITE, "font": FONT_HEADING},
])
add_rich_paragraphs(s, Inches(0.6), Inches(3.4), Inches(12), Inches(2), [
    {"text": "criterio antes que tecnología.", "size": 56, "bold": True, "color": MELI_YELLOW, "font": FONT_HEADING},
])

# Subtitle bullets
add_text(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
         "5 casos · 7 widgets en vivo · catálogo SOTA · defensa frente a stakeholders",
         size=16, italic=True, color=WHITE, font=FONT_BODY)

# Footer
add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.4),
         "andresislas2107@gmail.com · México · sandbox interactivo en URL pública",
         size=11, color=MELI_YELLOW_SOFT, font=FONT_BODY)


# ============================================================
# SLIDE 2 — AGENDA
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 2)
slide_title(s, "Agenda · 25 minutos", "Lo que vamos a recorrer juntos.",
            "Conclusión primero (Pyramid Principle Minto). Cada caso = 1 slide. Q&A al final.")

agenda_items = [
    ("01", "Resumen ejecutivo", "La conclusión + el porqué — 2 min"),
    ("02", "Caso 1 · Empaquetado", "Cómo decidir entre 5+ vendors — AHP scorecard — 3 min"),
    ("03", "Caso 2 · Bottleneck STP", "Por qué 1.6M items/mes no se alcanza — 3 min"),
    ("04", "Caso 3 · Picking robótico", "Productividad + palanca de mayor ROI — 3 min"),
    ("05", "Caso 4 · Plan data-driven", "Arquitectura + 3 queries clave — 2 min"),
    ("06", "Caso 5 · Roadmap Gantt", "28 sem a contrato + Real Options — 3 min"),
    ("07", "Bonus · Roadmap 18 meses", "Visión + métricas objetivo — 2 min"),
    ("08", "Defensa de gaps + Q&A", "Lo que sí, lo que no, lo que ofrezco — 5 min"),
]

y = Inches(2.3)
for num, title, sub in agenda_items:
    add_text(s, Inches(0.6), y, Inches(0.7), Inches(0.5),
             num, size=18, bold=True, color=MELI_BLUE, font=FONT_MONO)
    add_text(s, Inches(1.4), y, Inches(5), Inches(0.5),
             title, size=18, bold=True, color=INK, font=FONT_HEADING)
    add_text(s, Inches(6.5), y, Inches(6.5), Inches(0.5),
             sub, size=13, italic=True, color=MUTED, font=FONT_BODY)
    y += Inches(0.55)


# ============================================================
# SLIDE 3 — EXECUTIVE SUMMARY (SCQA)
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 3)
slide_title(s, "Resumen ejecutivo · SCQA",
            "Antes de elegir robots, pregunto si debe ser robot.",
            "Situation → Complication → Question → Answer (Minto, McKinsey).")

# 4 colored boxes side by side
box_y = Inches(2.4)
box_h = Inches(4.5)
box_w = Inches(3.05)
box_gap = Inches(0.15)
start_x = Inches(0.5)

boxes = [
    ("SITUATION", "Contexto", MELI_BLUE,
     "MELI escala agresivo en MX:\n\n• 3.2M sq ft Cuautitlán/Tepotzotlán 2023\n\n• 22 nuevos sitios logísticos 2024\n\n• Nuevo FC Hidalgo 100K m²\n\n• Cajamar Brasil ya con 334 robots STP\n\n• Acuerdo Agility 2025 (humanoides)"),
    ("COMPLICATION", "Tensión", DANGER,
     "Cada decisión de robotización:\n\n• CAPEX alto (USD 250K-3M+)\n\n• Payback incierto sin POC\n\n• Vendor lock-in latente\n\n• Pipeline de talento limitado\n\n• Risk de escalar mal lo que funciona en piloto"),
    ("QUESTION", "Pregunta clave", INFO,
     "¿Cómo evalúo cuándo, qué y dónde automatizar de forma defendible ante CEO/CFO/Ops?\n\nNo sólo:\n\"¿cuál robot?\"\n\nSino:\n\"¿debe ser robot — y bajo qué condiciones?\""),
    ("ANSWER", "Mi propuesta", SUCCESS,
     "5 lentes integradas:\n\n① Operativa (TOC bottleneck)\n② Técnica (fit + integración)\n③ Financiera (ROI + Real Options)\n④ Talento (pipeline universitario)\n⑤ Riesgo (vendor stability + lock-in)\n\nPor cada caso → recomendación gateada por POC."),
]

for i, (eyebrow, label, color, body) in enumerate(boxes):
    x = start_x + (box_w + box_gap) * i
    # Header band
    add_rect(s, x, box_y, box_w, Inches(0.55), fill_color=color)
    add_text(s, x + Inches(0.15), box_y + Inches(0.08), box_w - Inches(0.3), Inches(0.35),
             eyebrow, size=10, bold=True, color=WHITE, font=FONT_MONO)
    add_text(s, x + Inches(0.15), box_y + Inches(0.28), box_w - Inches(0.3), Inches(0.3),
             label, size=14, bold=True, color=WHITE, font=FONT_HEADING)
    # Body
    add_rect(s, x, box_y + Inches(0.55), box_w, box_h - Inches(0.55),
             fill_color=PAGE_BG, line_color=color, line_width=Pt(1.5))
    add_text(s, x + Inches(0.18), box_y + Inches(0.7),
             box_w - Inches(0.36), box_h - Inches(0.85),
             body, size=11, color=INK, font=FONT_BODY)


# ============================================================
# SLIDE 4 — CASO 1 EMPAQUETADO
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 4)
slide_title(s, "Caso 1 · Empaquetado",
            "Packsize gana con pesos default; Sealed Air si subes Soporte MX.",
            "AHP scorecard 6 criterios × 5 vendors, rúbrica con anchors 0/25/50/75/100.")

# Left: ranking
add_text(s, Inches(0.6), Inches(2.4), Inches(6), Inches(0.4),
         "Ranking AHP — pesos default Fit 25% · Sop 20% · Tal 10% · Eco 25% · Rie 15% · Sus 5%",
         size=11, bold=True, color=MELI_BLUE, font=FONT_BODY)

ranking = [
    ("#1 Packsize X-Series (auto-box)", 83, MELI_YELLOW),
    ("#2 Sparck CMC Pakto (auto-wrap)", 80, MELI_YELLOW_DARK),
    ("#3 Sealed Air AutoBag 850HB (bagging)", 72, MELI_BLUE),
    ("#4 Quadient CVP Impack (auto-box)", 70, INFO),
    ("#5 Status quo manual + Smurfit", 69, MUTED),
]
y = Inches(2.85)
for label, score, color in ranking:
    add_rect(s, Inches(0.6), y, Inches(0.15), Inches(0.45), fill_color=color)
    add_text(s, Inches(0.9), y, Inches(4.8), Inches(0.45),
             label, size=12, bold=True, color=INK, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    # Score bar background
    add_rect(s, Inches(5.7), y + Inches(0.1), Inches(1.2), Inches(0.25),
             fill_color=RGBColor(0xEF, 0xEF, 0xF5))
    # Score bar fill
    bar_w = Inches(1.2 * score / 100)
    add_rect(s, Inches(5.7), y + Inches(0.1), bar_w, Inches(0.25),
             fill_color=color)
    add_text(s, Inches(6.95), y, Inches(0.6), Inches(0.45),
             f"{score}", size=14, bold=True, color=MELI_BLUE, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)

# Right: insight box
insight_x = Inches(8.5)
insight_w = Inches(4.4)
add_rect(s, insight_x, Inches(2.4), insight_w, Inches(4.3),
         fill_color=MELI_YELLOW_SOFT, line_color=MELI_YELLOW_DARK, line_width=Pt(2))
add_text(s, insight_x + Inches(0.25), Inches(2.5),
         insight_w - Inches(0.5), Inches(0.4),
         "INSIGHT", size=11, bold=True, color=MELI_BLUE, font=FONT_MONO)

add_rich_paragraphs(s, insight_x + Inches(0.25), Inches(2.9),
                    insight_w - Inches(0.5), Inches(3.9), [
    {"text": "Packsize lidera con pesos balanceados",
     "size": 14, "bold": True, "color": MELI_BLUE},
    {"text": "porque combina presencia MX directa (case Autopartes Tracto verificable) + portafolio consolidado post-adquisición Sparck mayo 2025.",
     "size": 11, "color": INK, "space_before": 4},
    {"text": "Si subes peso de Soporte MX a 35%+, Sealed Air rebasa",
     "size": 14, "bold": True, "color": MELI_BLUE, "space_before": 14},
    {"text": "por sus 2 plantas MX (Toluca + Monterrey) y AutoBag 850HB de 700 órdenes/hr.",
     "size": 11, "color": INK, "space_before": 4},
    {"text": "Recomendación:",
     "size": 14, "bold": True, "color": MELI_BLUE, "space_before": 14},
    {"text": "RFI formal a Packsize y Sealed Air. Excluir Quadient hasta validar canal MX.",
     "size": 11, "color": INK, "space_before": 4},
])


# ============================================================
# SLIDE 5 — CASO 2 BOTTLENECK STP
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 5)
slide_title(s, "Caso 2 · Shelves-to-Person",
            "Put Away es el cuello severo. 1.6M items/mes no se alcanza con los parámetros dados.",
            "Theory of Constraints (Goldratt) + Little's Law aplicados a operación con utilización 10%.")

# Big KPI cards horizontal
kpi_y = Inches(2.4)
kpi_h = Inches(1.7)
kpi_w = Inches(2.95)
kpi_x = Inches(0.5)

kpis = [
    ("Put Away", "165K", "items/mes efectivos\n(10% utilización)", DANGER, "BOTTLENECK"),
    ("Picking", "1.54M", "items/mes efectivos\n(90% eficiencia)", INFO, ""),
    ("Sorter", "1.82M", "items/mes capacidad\n(sin factor)", SUCCESS, ""),
    ("Target", "1.6M", "items/mes meta\ngerencia", MELI_BLUE, "NO ALCANZA"),
]

for i, (label, value, sub, color, tag) in enumerate(kpis):
    x = kpi_x + (kpi_w + Inches(0.15)) * i
    add_rect(s, x, kpi_y, kpi_w, kpi_h, fill_color=PAGE_BG,
             line_color=color, line_width=Pt(2))
    add_text(s, x + Inches(0.15), kpi_y + Inches(0.1),
             kpi_w - Inches(0.3), Inches(0.3),
             label, size=11, bold=True, color=MUTED, font=FONT_MONO)
    add_text(s, x + Inches(0.15), kpi_y + Inches(0.4),
             kpi_w - Inches(0.3), Inches(0.7),
             value, size=44, bold=True, color=color, font=FONT_HEADING)
    add_text(s, x + Inches(0.15), kpi_y + Inches(1.15),
             kpi_w - Inches(0.3), Inches(0.4),
             sub, size=10, color=MUTED, font=FONT_BODY)
    if tag:
        add_rect(s, x + Inches(0.15), kpi_y + Inches(1.45),
                 Inches(1.6), Inches(0.22), fill_color=color)
        add_text(s, x + Inches(0.15), kpi_y + Inches(1.46),
                 Inches(1.6), Inches(0.22),
                 tag, size=9, bold=True, color=WHITE, font=FONT_MONO,
                 align=PP_ALIGN.CENTER)

# Bottom panel: palancas + Little's Law
panel_y = Inches(4.35)
panel_h = Inches(2.4)

# Palancas (left)
add_rect(s, Inches(0.5), panel_y, Inches(7.5), panel_h,
         fill_color=MELI_YELLOW_SOFT, line_color=MELI_YELLOW_DARK)
add_text(s, Inches(0.7), panel_y + Inches(0.1),
         Inches(7), Inches(0.35),
         "PALANCAS · sin agregar puestos (ordenadas por ROI marginal)",
         size=11, bold=True, color=MELI_BLUE, font=FONT_MONO)
add_rich_paragraphs(s, Inches(0.7), panel_y + Inches(0.55),
                    Inches(7.1), Inches(1.85), [
    {"text": "1. Put Away — subir utilización 10% → 97%", "size": 12, "bold": True, "color": INK},
    {"text": "    Root cause: ¿wave planning? ¿bursts inbound? Auditoría 2 sem.",
     "size": 10, "italic": True, "color": MUTED, "space_before": 2},
    {"text": "2. Picking — eficiencia 90% → 94%", "size": 12, "bold": True, "color": INK, "space_before": 8},
    {"text": "    Sólo 4 pp para cerrar gap residual. Slotting + batch picking.",
     "size": 10, "italic": True, "color": MUTED, "space_before": 2},
    {"text": "3. Subir PPF (caso 3) — sin tocar puestos", "size": 12, "bold": True, "color": INK, "space_before": 8},
    {"text": "    1.8 → 2.5 = +39% items/hr; cero CAPEX, sólo datos + algoritmo.",
     "size": 10, "italic": True, "color": MUTED, "space_before": 2},
])

# Little's Law (right)
add_rect(s, Inches(8.3), panel_y, Inches(4.6), panel_h,
         fill_color=MELI_BLUE)
add_text(s, Inches(8.5), panel_y + Inches(0.1),
         Inches(4.2), Inches(0.35),
         "LITTLE'S LAW · L = λW", size=11, bold=True, color=MELI_YELLOW,
         font=FONT_MONO)
add_text(s, Inches(8.5), panel_y + Inches(0.55),
         Inches(4.2), Inches(0.6),
         "WIP pick→sort", size=12, color=WHITE, font=FONT_BODY)
add_text(s, Inches(8.5), panel_y + Inches(0.95),
         Inches(4.2), Inches(0.9),
         "≈ 7,973", size=40, bold=True, color=MELI_YELLOW, font=FONT_HEADING)
add_text(s, Inches(8.5), panel_y + Inches(1.7),
         Inches(4.2), Inches(0.7),
         "items en buffer estable para SLA 3.5 hr.\n2,532 items/hr × 3.5 hr × factor 0.9.",
         size=10, color=WHITE, font=FONT_BODY)


# ============================================================
# SLIDE 6 — CASO 3 PICKING ROBÓTICO
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 6)
slide_title(s, "Caso 3 · Picking robótico",
            "301 items/hr/estación. PPF (units/rack) es la palanca de mayor ROI marginal.",
            "T_rack secuencial = OWP + OWS + R2R = 21.5 s. Combinando palancas: +76% sin CAPEX.")

# Formula box (left)
fb_x = Inches(0.5)
fb_w = Inches(5.5)
fb_h = Inches(4.8)
fb_y = Inches(2.4)
add_rect(s, fb_x, fb_y, fb_w, fb_h, fill_color=PAGE_BG,
         line_color=MELI_BLUE, line_width=Pt(2))
add_text(s, fb_x + Inches(0.2), fb_y + Inches(0.15),
         fb_w - Inches(0.4), Inches(0.4),
         "FÓRMULA", size=11, bold=True, color=MELI_BLUE, font=FONT_MONO)

# Big formula
add_text(s, fb_x + Inches(0.2), fb_y + Inches(0.6),
         fb_w - Inches(0.4), Inches(0.55),
         "T_rack = OWP + OWS + R2R",
         size=22, bold=True, color=INK, font=FONT_MONO)
add_text(s, fb_x + Inches(0.2), fb_y + Inches(1.15),
         fb_w - Inches(0.4), Inches(0.4),
         "         = 7 + 6 + 8.5 = 21.5 s",
         size=18, color=MUTED, font=FONT_MONO)

# Big KPIs
add_text(s, fb_x + Inches(0.2), fb_y + Inches(1.85),
         fb_w - Inches(0.4), Inches(0.35),
         "Racks/hr = 3,600 / 21.5 =",
         size=13, color=INK, font=FONT_BODY)
add_text(s, fb_x + Inches(0.2), fb_y + Inches(2.15),
         fb_w - Inches(0.4), Inches(0.7),
         "167.4 racks/hr",
         size=32, bold=True, color=MELI_BLUE, font=FONT_HEADING)

add_text(s, fb_x + Inches(0.2), fb_y + Inches(3.0),
         fb_w - Inches(0.4), Inches(0.35),
         "Items/hr = 167.4 × 1.8 PPF =",
         size=13, color=INK, font=FONT_BODY)
add_text(s, fb_x + Inches(0.2), fb_y + Inches(3.3),
         fb_w - Inches(0.4), Inches(0.7),
         "301.4 items/hr/estación",
         size=32, bold=True, color=SUCCESS, font=FONT_HEADING)

add_text(s, fb_x + Inches(0.2), fb_y + Inches(4.2),
         fb_w - Inches(0.4), Inches(0.5),
         "(Por qué secuencial: enunciado dice \"una vez confirmado, robots inician movimiento\".)",
         size=10, italic=True, color=MUTED, font=FONT_BODY)

# Sensitivity panel (right)
sb_x = Inches(6.3)
sb_w = Inches(6.6)
sb_h = Inches(4.8)
sb_y = Inches(2.4)
add_rect(s, sb_x, sb_y, sb_w, sb_h, fill_color=MELI_BLUE)
add_text(s, sb_x + Inches(0.2), sb_y + Inches(0.15),
         sb_w - Inches(0.4), Inches(0.4),
         "SENSIBILIDAD POR PALANCA · Δ vs baseline",
         size=11, bold=True, color=MELI_YELLOW, font=FONT_MONO)

sens = [
    ("PPF 1.8 → 2.5", "+39%", "418 items/hr", "Slotting algoritmo — CERO CAPEX"),
    ("R2R 8.5 → 6.0 s", "+13%", "340 items/hr", "Más robots, paths optimizados"),
    ("OWP 7 → 5 s", "+10%", "332 items/hr", "Light-directed picking, mejor UI"),
    ("OWS 6 → 4 s", "+10%", "332 items/hr", "Ergonomía + ubicación de totes"),
    ("Combinado", "+76%", "530 items/hr", "Slotting + R2R + OWP simultáneo"),
]
y = sb_y + Inches(0.65)
for desc, pct, val, note in sens:
    is_combined = "Combinado" in desc
    bg = MELI_YELLOW if is_combined else MELI_BLUE_DARK
    color_pct = MELI_BLUE_DARK if is_combined else MELI_YELLOW
    color_text = MELI_BLUE_DARK if is_combined else WHITE
    add_rect(s, sb_x + Inches(0.2), y, sb_w - Inches(0.4), Inches(0.7),
             fill_color=bg)
    add_text(s, sb_x + Inches(0.35), y + Inches(0.06),
             Inches(1.9), Inches(0.6),
             desc, size=12, bold=True, color=color_text, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, sb_x + Inches(2.3), y + Inches(0.06),
             Inches(0.9), Inches(0.6),
             pct, size=18, bold=True, color=color_pct, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, sb_x + Inches(3.2), y + Inches(0.06),
             Inches(1.3), Inches(0.6),
             val, size=11, color=color_text, font=FONT_MONO,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, sb_x + Inches(4.5), y + Inches(0.06),
             Inches(2.0), Inches(0.6),
             note, size=9, italic=True, color=color_text, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.78)


# ============================================================
# SLIDE 7 — CASO 4 PLAN DATA
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 7)
slide_title(s, "Caso 4 · Plan data-driven",
            "Stack asumido + 3 queries SQL + dashboard con insight accionable.",
            "Más allá del dashboard pasivo: cada vista termina en una recomendación concreta.")

# Architecture (left)
arch_x = Inches(0.5)
arch_w = Inches(5.5)
arch_y = Inches(2.4)
add_rect(s, arch_x, arch_y, arch_w, Inches(4.3),
         fill_color=PAGE_BG, line_color=MELI_BLUE, line_width=Pt(2))
add_text(s, arch_x + Inches(0.2), arch_y + Inches(0.15),
         arch_w - Inches(0.4), Inches(0.4),
         "ARQUITECTURA · stack asumido (S001)",
         size=11, bold=True, color=MELI_BLUE, font=FONT_MONO)

arch_steps = [
    ("WMS", "custom MELI / SAP EWM"),
    ("↓ Kafka events", "streaming en tiempo real"),
    ("Data Lake", "S3 / GCS + Glue / Athena"),
    ("↓ ETL incremental", "transformación schema-on-read"),
    ("Data Warehouse", "Redshift / Snowflake / BigQuery"),
    ("↓ semantic layer", "agregaciones reutilizables"),
    ("BI tools", "Looker / Tableau / Streamlit ad-hoc"),
]
y = arch_y + Inches(0.7)
for label, sub in arch_steps:
    is_arrow = label.startswith("↓")
    color = MUTED if is_arrow else MELI_BLUE
    size = 11 if is_arrow else 14
    bold = not is_arrow
    add_text(s, arch_x + Inches(0.3), y, arch_w - Inches(0.6), Inches(0.3),
             label, size=size, bold=bold, color=color, font=FONT_MONO)
    add_text(s, arch_x + Inches(2.3), y, arch_w - Inches(2.6), Inches(0.3),
             sub, size=9, italic=True, color=MUTED, font=FONT_BODY)
    y += Inches(0.45)

# Queries (right)
q_x = Inches(6.3)
q_w = Inches(6.6)
q_y = Inches(2.4)
add_rect(s, q_x, q_y, q_w, Inches(4.3),
         fill_color=MELI_BLUE)
add_text(s, q_x + Inches(0.2), q_y + Inches(0.15),
         q_w - Inches(0.4), Inches(0.4),
         "3 QUERIES + INSIGHT ACCIONABLE",
         size=11, bold=True, color=MELI_YELLOW, font=FONT_MONO)

queries = [
    ("① Curvas semanales entrada/salida", "GROUP BY fc, week, event_type"),
    ("② Curva intradiaria + peaks", "EXTRACT hour + PERCENTILE_CONT(0.95)"),
    ("③ Mix SKU + dim-weight", "JOIN sku_master + WIDTH_BUCKET dim_weight"),
]
y = q_y + Inches(0.6)
for q, snippet in queries:
    add_text(s, q_x + Inches(0.3), y,
             q_w - Inches(0.6), Inches(0.35),
             q, size=13, bold=True, color=WHITE, font=FONT_BODY)
    add_text(s, q_x + Inches(0.3), y + Inches(0.35),
             q_w - Inches(0.6), Inches(0.3),
             snippet, size=10, color=MELI_YELLOW, font=FONT_MONO)
    y += Inches(0.8)

add_rect(s, q_x + Inches(0.2), y + Inches(0.1),
         q_w - Inches(0.4), Inches(1.1),
         fill_color=MELI_YELLOW)
add_text(s, q_x + Inches(0.35), y + Inches(0.18),
         q_w - Inches(0.7), Inches(0.3),
         "INSIGHT SLIDE PARA EL LÍDER",
         size=10, bold=True, color=MELI_BLUE_DARK, font=FONT_MONO)
add_text(s, q_x + Inches(0.35), y + Inches(0.5),
         q_w - Inches(0.7), Inches(0.6),
         "\"FC-X: utilización 60% pero throughput plano vs FC-Y.\nHipótesis: bottleneck en sortation post-pick.\nRecomendación: profundizar antes de expandir Y a más FCs.\"",
         size=11, italic=True, color=MELI_BLUE_DARK, font=FONT_BODY)


# ============================================================
# SLIDE 8 — CASO 5 ROADMAP GANTT
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 8)
slide_title(s, "Caso 5 · Roadmap feasibility-a-licitación",
            "28 sem a contrato firmado · 46 sem a steady state · gate Fase 7 = Real Options.",
            "CPM + PERT Beta-PERT (Vose 2008) + Real Options binomial (Cox-Ross-Rubinstein 1979).")

# Timeline visualization
tl_x = Inches(0.5)
tl_y = Inches(2.5)
tl_w = Inches(12.3)
tl_h = Inches(0.55)

# Phases as horizontal bars
phases = [
    ("Discovery + RFI", 0, 9, MELI_BLUE),
    ("POC", 9, 15, SUCCESS),
    ("Business case + gate", 15, 20, DANGER),  # gate fase 7 marker
    ("RFP + licitación + contrato", 20, 28, MELI_BLUE),
    ("Fabricación (dado)", 28, 40, MELI_YELLOW),
    ("Site prep ∥", 28, 32, MELI_YELLOW_DARK),
    ("Install + commissioning", 40, 44, MELI_BLUE),
    ("Ramp-up", 44, 50, SUCCESS),
]

# Background ruler
add_rect(s, tl_x, tl_y, tl_w, Inches(0.04),
         fill_color=RGBColor(0xD8, 0xD8, 0xE3))
for w in range(0, 51, 10):
    x_tick = tl_x + Inches(tl_w.inches * w / 50)
    add_rect(s, x_tick, tl_y, Emu(15000), Inches(0.08),
             fill_color=MUTED)
    add_text(s, x_tick - Inches(0.3), tl_y - Inches(0.35),
             Inches(0.6), Inches(0.25),
             f"sem {w}", size=8, color=MUTED, font=FONT_MONO,
             align=PP_ALIGN.CENTER)

# Phase bars
bar_y = tl_y + Inches(0.3)
bar_h_each = Inches(0.35)
prev_end = -1
row = 0
for label, start, end, color in phases:
    # if overlap with prev: bump to next row
    if start < prev_end and "site prep" not in label.lower():
        row += 1
    else:
        if "site prep" not in label.lower():
            row = 0
    sx = tl_x + Inches(tl_w.inches * start / 50)
    sw = Inches(tl_w.inches * (end - start) / 50)
    by = bar_y + Inches(row * 0.45)
    add_rect(s, sx, by, sw, bar_h_each, fill_color=color)
    add_text(s, sx + Inches(0.05), by, sw - Inches(0.1), bar_h_each,
             label, size=9, bold=True, color=WHITE, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    prev_end = end

# Gate marker (vertical red line at week 17, fase 7 approval)
gate_x = tl_x + Inches(tl_w.inches * 17 / 50)
add_rect(s, gate_x, tl_y - Inches(0.4), Inches(0.03), Inches(2.3),
         fill_color=DANGER)
add_text(s, gate_x + Inches(0.1), tl_y - Inches(0.4),
         Inches(2), Inches(0.3),
         "← Gate Fase 7 (Real Options)",
         size=9, bold=True, color=DANGER, font=FONT_BODY)

# Bottom: 3 conclusion cards
cc_y = Inches(5.0)
cc_h = Inches(1.7)
cc_w = Inches(4.0)
cc_x = Inches(0.5)
cc_gap = Inches(0.15)

conclusions = [
    ("28 SEMANAS", "a contrato firmado",
     "Pre-fabricación = 28 sem (~7 meses).\nIncluye discovery + POC + business case + licitación + negociación.",
     MELI_BLUE),
    ("46 SEMANAS", "a steady state",
     "Total = 46 sem (~11 meses).\nIncluye fabricación 12 sem + ramp-up 6 sem.\nSitio en paralelo a fab ahorra 4 sem.",
     SUCCESS),
    ("REAL OPTIONS", "gate Fase 7 = put option",
     "Si POC falla, no se ejerce CAPEX de fab.\nValor de la flexibilidad ≠ NPV \"punto\".\nWidget en deck web cuantifica.",
     DANGER),
]
for i, (kpi, sub, body, color) in enumerate(conclusions):
    x = cc_x + (cc_w + cc_gap) * i
    add_rect(s, x, cc_y, cc_w, cc_h, fill_color=PAGE_BG,
             line_color=color, line_width=Pt(2))
    add_text(s, x + Inches(0.2), cc_y + Inches(0.15),
             cc_w - Inches(0.4), Inches(0.4),
             kpi, size=18, bold=True, color=color, font=FONT_HEADING)
    add_text(s, x + Inches(0.2), cc_y + Inches(0.55),
             cc_w - Inches(0.4), Inches(0.3),
             sub, size=11, italic=True, color=MUTED, font=FONT_BODY)
    add_text(s, x + Inches(0.2), cc_y + Inches(0.85),
             cc_w - Inches(0.4), cc_h - Inches(1),
             body, size=10, color=INK, font=FONT_BODY)


# ============================================================
# SLIDE 9 — BONUS ROADMAP 18M
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 9)
slide_title(s, "Bonus · Roadmap 18 meses MELI MX Innovación",
            "Cómo abordaría la innovación operativa si tengo el rol.",
            "Lectura opcional — visión personal aplicada, no parte del brief.")

# 4 timeline blocks horizontal
tl_y = Inches(2.4)
tl_h = Inches(2.6)
tl_w = Inches(3.05)
tl_x = Inches(0.5)
tl_gap = Inches(0.15)

phases_18m = [
    ("MES 0-3", "Assessment", MELI_BLUE,
     "• 3 CEDIS principales\n• Bottleneck map\n• Score fit por tech\n• Gate: priorizar 2-3 iniciativas",
     "Comité prioriza iniciativas para POC"),
    ("MES 3-6", "Pilotos rápidos", SUCCESS,
     "• Auto-box packaging en 1 FC\n• Algoritmo slotting (cero CAPEX)\n• POC 6 sem c/u\n• Gate: ROI validado",
     "Real Options aplicado en cada gate"),
    ("MES 6-12", "Rollout + nuevas", MELI_YELLOW_DARK,
     "• Escalado del piloto exitoso\n• POC STP/AMR en FC maduro\n• Sorter upgrade si aplica\n• Visita técnica GEODIS-Locus",
     "Cuautitlán como benchmark local"),
    ("MES 12-18", "Arquitectura + talento", DANGER,
     "• Unified WMS ↔ robots (ISA-95 + OPC UA)\n• Predictive maintenance fleet\n• Programa universitario 3-4 univ.\n• Lab sponsored conjunto",
     "Pipeline talento sostenible"),
]

for i, (period, title, color, body, footer) in enumerate(phases_18m):
    x = tl_x + (tl_w + tl_gap) * i
    add_rect(s, x, tl_y, tl_w, Inches(0.5), fill_color=color)
    add_text(s, x + Inches(0.15), tl_y + Inches(0.05),
             tl_w - Inches(0.3), Inches(0.2),
             period, size=10, bold=True, color=WHITE, font=FONT_MONO)
    add_text(s, x + Inches(0.15), tl_y + Inches(0.22),
             tl_w - Inches(0.3), Inches(0.3),
             title, size=14, bold=True, color=WHITE, font=FONT_HEADING)
    add_rect(s, x, tl_y + Inches(0.5), tl_w, tl_h - Inches(0.5),
             fill_color=PAGE_BG, line_color=color, line_width=Pt(1.5))
    add_text(s, x + Inches(0.18), tl_y + Inches(0.65),
             tl_w - Inches(0.36), Inches(1.5),
             body, size=10, color=INK, font=FONT_BODY)
    add_text(s, x + Inches(0.18), tl_y + Inches(2.15),
             tl_w - Inches(0.36), Inches(0.4),
             footer, size=9, italic=True, color=color, font=FONT_BODY)

# Bottom metrics
m_y = Inches(5.3)
m_h = Inches(1.4)

metrics_18m = [
    ("↓ 12–18%", "cost-per-unit", SUCCESS, "auto-box + slotting + STP combinados"),
    ("↑ 35–50%", "picks/hr/estación", SUCCESS, "slotting (~30%) + STP (~20% adic.)"),
    ("92% → 97%+", "uptime fleet target", SUCCESS, "predictive maintenance + refacciones"),
    ("18–24 mo", "payback agregado", MELI_BLUE, "AMR <24mo + packaging 1-2yr (Mordor)"),
]
mw = Inches(3.05)
mx = Inches(0.5)
mgap = Inches(0.15)
for i, (val, label, color, sub) in enumerate(metrics_18m):
    x = mx + (mw + mgap) * i
    add_rect(s, x, m_y, mw, m_h, fill_color=MELI_BLUE_DARK)
    add_text(s, x + Inches(0.15), m_y + Inches(0.1),
             mw - Inches(0.3), Inches(0.5),
             val, size=22, bold=True, color=color, font=FONT_HEADING)
    add_text(s, x + Inches(0.15), m_y + Inches(0.65),
             mw - Inches(0.3), Inches(0.3),
             label, size=11, color=WHITE, font=FONT_BODY)
    add_text(s, x + Inches(0.15), m_y + Inches(0.95),
             mw - Inches(0.3), Inches(0.4),
             sub, size=9, italic=True, color=MELI_YELLOW_SOFT, font=FONT_BODY)


# ============================================================
# SLIDE 10 — DEFENSA DE GAPS
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 10)
slide_title(s, "Defensa de gaps",
            "Lo que NO tengo · cómo lo cubro · qué ofrezco a cambio.",
            "Declarar los gaps abre puerta a discusión honesta. Inflar los cierra.")

gaps = [
    ("Fleet máx 4 vehículos",
     "MELI maneja 10–100× esa escala.",
     "Adyacencia: criterio de escalado progresivo (asesoría cliente CEDIS-regulado) + arquitectura OPC UA → WMS → tareas robot. Lo que necesito aprender: dimensionamiento fleet software a escala (MAPF), social dynamics de 100+ operadores."),
    ("Primer rol e-commerce directo",
     "Mis 7 años son en planta industrial pesada con JIT.",
     "Adyacencia: dinámicas TOC, OEE, mantenibilidad, wave-planning son las mismas. Lo que necesito aprender: SKU mix B2C high-variety, cadencia Buen Fin, last-mile."),
    ("Portugués no fluido",
     "Mi inglés es C1, mi español nativo, portugués técnico ~70% inteligible pero no produzco fluido.",
     "Plan: ES + EN cubren 95% de situaciones con Brasil. Si el rol requiere PT activo, lo aprendo (~3 meses con clases + uso diario)."),
]

y = Inches(2.4)
gap_h = Inches(1.45)
for label, what, mitigation in gaps:
    add_rect(s, Inches(0.5), y, Inches(0.15), gap_h, fill_color=MELI_YELLOW_DARK)
    add_rect(s, Inches(0.65), y, Inches(12.2), gap_h,
             fill_color=PAGE_BG, line_color=RGBColor(0xD8, 0xD8, 0xE3))
    add_text(s, Inches(0.85), y + Inches(0.1),
             Inches(4), Inches(0.4),
             label, size=15, bold=True, color=MELI_BLUE, font=FONT_HEADING)
    add_text(s, Inches(0.85), y + Inches(0.5),
             Inches(11.5), Inches(0.35),
             what, size=11, italic=True, color=MUTED, font=FONT_BODY)
    add_text(s, Inches(0.85), y + Inches(0.85),
             Inches(11.5), gap_h - Inches(0.95),
             mitigation, size=11, color=INK, font=FONT_BODY)
    y += gap_h + Inches(0.15)


# ============================================================
# SLIDE 11 — CIERRE + URL WEB
# ============================================================

s = prs.slides.add_slide(BLANK)
# Background MELI blue
add_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill_color=MELI_BLUE)
add_rect(s, Emu(0), Emu(0), Inches(2), Inches(0.5), fill_color=MELI_YELLOW)

# Title
add_text(s, Inches(0.6), Inches(1.0),
         Inches(12), Inches(0.4),
         "CIERRE",
         size=12, bold=True, color=MELI_YELLOW, font=FONT_MONO)
add_text(s, Inches(0.6), Inches(1.5),
         Inches(12), Inches(1.0),
         "Gracias.",
         size=52, bold=True, color=WHITE, font=FONT_HEADING)
add_text(s, Inches(0.6), Inches(2.5),
         Inches(12), Inches(0.5),
         "Lo importante de este case no fue \"qué robot\".",
         size=18, color=MELI_YELLOW, font=FONT_BODY)
add_text(s, Inches(0.6), Inches(3.0),
         Inches(12), Inches(0.5),
         "Fue cuándo + por qué + cómo defenderlo. Esa es la conversación que quiero abrir con ustedes.",
         size=18, color=WHITE, font=FONT_BODY)

# URL box
url_y = Inches(4.2)
add_rect(s, Inches(0.6), url_y, Inches(12), Inches(1.7),
         fill_color=MELI_BLUE_DARK, line_color=MELI_YELLOW, line_width=Pt(2))
add_text(s, Inches(0.85), url_y + Inches(0.15),
         Inches(11.5), Inches(0.3),
         "EXPLOREN EL SANDBOX INTERACTIVO",
         size=10, bold=True, color=MELI_YELLOW, font=FONT_MONO)
add_text(s, Inches(0.85), url_y + Inches(0.45),
         Inches(11.5), Inches(0.5),
         "https://andresislas99.github.io/meli-innovation-business-case/",
         size=18, bold=True, color=WHITE, font=FONT_MONO)
add_text(s, Inches(0.85), url_y + Inches(0.95),
         Inches(11.5), Inches(0.65),
         "9 slides · 7 widgets Python en vivo · rúbrica AHP con evidencia por vendor · simulador DES · Real Options binomial · Sobol indices · roadmap 18 meses · guía de estudio + cheatsheet · catálogo SOTA con papers canónicos.",
         size=12, color=WHITE, font=FONT_BODY)

# Contact
add_text(s, Inches(0.6), Inches(6.3),
         Inches(12), Inches(0.4),
         "Andrés Islas Bravo · andresislas2107@gmail.com · México",
         size=14, italic=True, color=MELI_YELLOW, font=FONT_BODY,
         align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — GLOSARIO RÁPIDO (Top 30)
# ============================================================

s = prs.slides.add_slide(BLANK)
add_brand_strip(s, 12)
slide_title(s, "Glosario · siglas que pueden aparecer",
            "Para referencia rápida durante Q&A — glosario extendido en exports/glosario.pdf.",
            "Cada sigla con significado pleno + dónde aparece en este deck.")

glosario_left = [
    ("TOC", "Theory of Constraints", "Caso 2"),
    ("AHP", "Analytic Hierarchy Process", "Caso 1"),
    ("OEE", "Overall Equipment Effectiveness", "Caso 2"),
    ("CPM", "Critical Path Method", "Caso 5"),
    ("PERT", "Program Evaluation/Review Technique", "Caso 5"),
    ("WIP", "Work-in-Process", "Caso 2"),
    ("CAPEX", "Capital Expenditure (inversión)", "Todos"),
    ("OPEX", "Operating Expenditure (recurrente)", "Todos"),
    ("NPV", "Net Present Value (valor presente)", "Caso 5"),
    ("ROI", "Return on Investment", "Todos"),
    ("TCO", "Total Cost of Ownership 5y", "Caso 1"),
    ("RFI / RFP", "Request for Info/Proposal", "Caso 1, 5"),
    ("SLA", "Service Level Agreement", "Caso 1"),
    ("MTTR / MTBF", "Mean Time To Repair / Between Failures", "Caso 1"),
    ("WMS / ERP / MES", "Warehouse / Enterprise / Manuf. Execution System", "Caso 1, 4"),
]
glosario_right = [
    ("FC / CEDIS", "Fulfillment Center / Centro Distribución", "Todos"),
    ("STP / AMR / AGV", "Shelves-to-Person / Autonomous Mobile / Auto-Guided", "Caso 2, 3"),
    ("OWP / OWS / R2R", "Operator Walk Pick / Stow + Rack-to-Rack", "Caso 3"),
    ("PPF", "Pieces Per Fetch (items por rack)", "Caso 3"),
    ("DIM weight", "Peso volumétrico (UPS/FedEx)", "Caso 1"),
    ("MAPF / PIBT", "Multi-Agent Path Finding", "SOTA"),
    ("CP-SAT", "Constraint Programming SAT (OR-Tools)", "SOTA"),
    ("DES", "Discrete Event Simulation", "Caso 2 SOTA"),
    ("MDP / RL / MARL", "Markov / Reinforcement / Multi-Agent RL", "SOTA"),
    ("OPC UA", "Industrial protocol (machine-to-WMS)", "Caso 1, 5"),
    ("ISA-95", "Enterprise-control integration standard", "Caso 1, 5"),
    ("RCM", "Reliability-Centered Maintenance", "Caso 1"),
    ("ABC / XYZ", "Pareto vol/value + variability", "Caso 4"),
    ("S₁ / S_T", "Sobol indices first/total order", "Caso 5 SOTA"),
    ("u / d / p", "Up/down/prob en binomial tree", "Caso 5 SOTA"),
]

col_y = Inches(2.4)
col_w = Inches(6.1)
row_h = Inches(0.3)

for col_idx, col_data in enumerate([glosario_left, glosario_right]):
    col_x = Inches(0.5) + col_w * col_idx + Inches(0.15) * col_idx
    y = col_y
    for sigla, meaning, where in col_data:
        add_text(s, col_x, y, Inches(1.5), row_h,
                 sigla, size=11, bold=True, color=MELI_BLUE, font=FONT_MONO)
        add_text(s, col_x + Inches(1.5), y, Inches(3.6), row_h,
                 meaning, size=10, color=INK, font=FONT_BODY)
        add_text(s, col_x + Inches(5.1), y, Inches(1.0), row_h,
                 where, size=9, italic=True, color=MUTED, font=FONT_BODY,
                 align=PP_ALIGN.RIGHT)
        y += row_h


# ============================================================
# SAVE
# ============================================================

prs.save(str(OUTPUT))
print(f"✅ {OUTPUT} ({OUTPUT.stat().st_size // 1024} KB, {len(prs.slides)} slides)")
